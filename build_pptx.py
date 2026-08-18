#!/usr/bin/env python3
"""Generate threatlint executive overview PPTX — concise 6-slide deck."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

NAVY   = RGBColor(0x0D, 0x1B, 0x2A)
AMBER  = RGBColor(0xD9, 0x77, 0x06)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHT = RGBColor(0xF4, 0xF6, 0xF9)
GREY   = RGBColor(0x6B, 0x7A, 0x8D)
INK    = RGBColor(0x1A, 0x23, 0x30)
GREEN  = RGBColor(0x2E, 0xCC, 0x71)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]


def rect(slide, x, y, w, h, rgb):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.line.fill.background()
    s.fill.solid()
    s.fill.fore_color.rgb = rgb
    return s


def txt(slide, text, x, y, w, h, size=16, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def light_slide(slide, title):
    rect(slide, 0, 0, W, H, OFFWHT)
    rect(slide, 0, 0, W, Inches(1.05), NAVY)
    rect(slide, 0, Inches(1.05), Inches(0.06), H - Inches(1.05), AMBER)
    txt(slide, title, Inches(0.22), Inches(0.12), Inches(12.5), Inches(0.82),
        size=28, bold=True, color=WHITE)


# ── Slide 1: Title ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
rect(s, 0, 0, W, H, NAVY)
rect(s, 0, H - Inches(0.1), W, Inches(0.1), AMBER)
rect(s, 0, Inches(3.15), W, Inches(0.05), AMBER)
txt(s, "threatlint", Inches(1), Inches(1.2), Inches(11.33), Inches(1.8),
    size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "AI-powered application security — everywhere you work",
    Inches(1), Inches(3.35), Inches(11.33), Inches(0.65),
    size=22, color=AMBER, align=PP_ALIGN.CENTER)
txt(s, "14 agents · Claude · OpenAI · Local LLMs · Claude Code · Copilot Chat · Teams · CI/CD",
    Inches(1), Inches(4.1), Inches(11.33), Inches(0.5),
    size=15, color=GREY, align=PP_ALIGN.CENTER)

# ── Slide 2: Problem → Solution ───────────────────────────────────────────────
s = prs.slides.add_slide(blank)
light_slide(s, "The Problem — and the Answer")

rect(s, Inches(0.3), Inches(1.2), Inches(6.0), Inches(5.85), WHITE)
rect(s, Inches(0.3), Inches(1.2), Inches(6.0), Inches(0.44), RGBColor(0xE7,0x4C,0x3C))
txt(s, "Security today", Inches(0.45), Inches(1.24), Inches(5.7), Inches(0.36),
    size=15, bold=True, color=WHITE)
problems = [
    "Manual reviews are a bottleneck",
    "Scanners generate noise, not context",
    "Threat models go stale immediately",
    "No single tool covers the full stack",
    "Compliance mapping is manual drudgery",
]
for i, p in enumerate(problems):
    txt(s, "✕  " + p, Inches(0.5), Inches(1.8 + 0.92 * i),
        Inches(5.6), Inches(0.75), size=15, color=INK)

rect(s, Inches(6.7), Inches(1.2), Inches(6.33), Inches(5.85), WHITE)
rect(s, Inches(6.7), Inches(1.2), Inches(6.33), Inches(0.44), RGBColor(0x2E,0xCC,0x71))
txt(s, "threatlint", Inches(6.85), Inches(1.24), Inches(6.03), Inches(0.36),
    size=15, bold=True, color=WHITE)
answers = [
    "14 agents — one per security domain",
    "Every finding cites file:line from real code",
    "Works on any repo, on demand or in CI",
    "Read-only: never modifies the workspace",
    "Structured output: executive + technical tiers",
]
for i, a in enumerate(answers):
    txt(s, "✓  " + a, Inches(6.85), Inches(1.8 + 0.92 * i),
        Inches(6.03), Inches(0.75), size=15, color=INK)

# ── Slide 3: The 14 Agents ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
light_slide(s, "14 Specialised Security Agents")

agents = [
    ("Threat Modeler",    "STRIDE · ATT&CK · DREAD",     NAVY),
    ("Code Reviewer",     "Diff / PR security review",   NAVY),
    ("Dependency Auditor","CVEs · Confusion · Hooks",     NAVY),
    ("Secrets Scanner",   "Keys · Tokens · Entropy",      NAVY),
    ("IaC Reviewer",      "Terraform · K8s · Dockerfile", NAVY),
    ("CI/CD Auditor",     "Injection · Pinning · Perms",  NAVY),
    ("API Security",      "OWASP API Top 10 (2023)",      NAVY),
    ("Auth Reviewer",     "OAuth · JWT · RBAC · MFA",     NAVY),
    ("FP Reviewer",       "Triage + Semgrep tuning",      NAVY),
    ("Compliance",        "ASVS · PCI · HIPAA · SOC 2",   AMBER),
    ("Attack Tree",       "AND/OR · Bypass · Leaf rank",  AMBER),
    ("Red Team",          "5 ATT&CK kill chains + IoCs",  AMBER),
    ("Threat Delta",      "RESOLVED / REGRESSED / NEW",   AMBER),
    ("Verify Fix",        "REMEDIATED / STILL PRESENT",   AMBER),
]

CW, CH = Inches(1.82), Inches(2.55)
for i, (name, desc, bg) in enumerate(agents):
    row, col = divmod(i, 7)
    x = Inches(0.24 + 1.84 * col)
    y = Inches(1.2 + 2.68 * row)
    rect(s, x, y, CW, CH, WHITE)
    rect(s, x, y, CW, Inches(0.44), bg)
    txt(s, name, x + Inches(0.1), y + Inches(0.06), CW - Inches(0.2),
        Inches(0.34), size=12, bold=True, color=WHITE)
    txt(s, desc, x + Inches(0.1), y + Inches(0.55), CW - Inches(0.2),
        Inches(1.8), size=11, color=INK)

# ── Slide 4: Works Everywhere ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
light_slide(s, "Works Everywhere")

platforms = [
    ("Claude Code",        NAVY,                          ["14 agents via @ picker", "24 slash commands", "Word doc output"]),
    ("Copilot Chat",       RGBColor(0x1F,0x88,0x3D),     ["Commit .github/agents/", "No install — reload VS Code", "All 14 agents in @ picker"]),
    ("Microsoft Teams",    RGBColor(0x62,0x64,0xA7),     ["Node.js Bot Framework bot", "14 commands + /help", "Claude API powered"]),
    ("GitHub Actions",     RGBColor(0x23,0x8B,0x9D),     ["PR review on every non-draft PR", "Weekly threat model schedule", "SARIF → Code Scanning"]),
    ("OpenAI / GPT",       RGBColor(0x41,0x29,0x91),     ["gpt-5.6 family", "--provider openai flag", "Same agents, same output"]),
    ("LM Studio (local)",  RGBColor(0x6D,0x28,0xD9),     ["100% on-device, no API key", "Air-gap / confidential repos", "Any 7B+ instruction model"]),
]

for i, (name, color, bullets) in enumerate(platforms):
    row, col = divmod(i, 3)
    x = Inches(0.25 + 4.28 * col)
    y = Inches(1.2 + 3.0 * row)
    rect(s, x, y, Inches(4.1), Inches(2.75), WHITE)
    rect(s, x, y, Inches(4.1), Inches(0.46), color)
    rect(s, x, y, Inches(0.06), Inches(2.75), color)
    txt(s, name, x + Inches(0.14), y + Inches(0.06), Inches(3.8), Inches(0.34),
        size=15, bold=True, color=WHITE)
    for j, b in enumerate(bullets):
        txt(s, "• " + b, x + Inches(0.18), y + Inches(0.62 + 0.68 * j),
            Inches(3.74), Inches(0.58), size=13, color=INK)

# ── Slide 5: How to Use It ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
light_slide(s, "How to Use It")

panels = [
    ("Claude Code", NAVY,
"""/threat-model
/threat-model-deep
/security-review main..feature
/compliance-check
/attack-tree src/payments
/red-team
/verify-fix CR-042"""),
    ("Copilot Chat / Teams", RGBColor(0x1F,0x88,0x3D),
"""@AppSec Threat Modeler
  threat model the payments service

@AppSec Red Team
  generate adversarial scenarios

@threatlint /compliance-check
@threatlint /help"""),
    ("GitHub Actions (auto)", RGBColor(0x23,0x8B,0x9D),
"""# Every non-draft PR:
appsec-pr-review.yml → PR comment + SARIF

# Weekly + IaC push to main:
appsec-scheduled.yml → GitHub Issues

# Secrets: ANTHROPIC_API_KEY
# or OPENAI_API_KEY
# or GITHUB_TOKEN (free, no key)"""),
]

for i, (title, color, code) in enumerate(panels):
    x = Inches(0.25 + 4.28 * i)
    rect(s, x, Inches(1.2), Inches(4.1), Inches(5.9), WHITE)
    rect(s, x, Inches(1.2), Inches(4.1), Inches(0.44), color)
    txt(s, title, x + Inches(0.12), Inches(1.25), Inches(3.86), Inches(0.34),
        size=14, bold=True, color=WHITE)
    rect(s, x + Inches(0.1), Inches(1.74), Inches(3.9), Inches(5.26), NAVY)
    txt(s, code, x + Inches(0.22), Inches(1.84), Inches(3.68), Inches(5.06),
        size=11, color=GREEN)

# ── Slide 6: Get Started ──────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
rect(s, 0, 0, W, H, NAVY)
rect(s, 0, H - Inches(0.1), W, Inches(0.1), AMBER)

txt(s, "Get Started in 3 Steps", Inches(1), Inches(0.55), Inches(11.33), Inches(0.75),
    size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

steps = [
    ("1  Copy agents",
     "cp threatlint/.claude/agents/*.md\n   ~/.claude/agents/\n\ncp threatlint/.claude/commands/*.md\n   ~/.claude/commands/"),
    ("2  Open any repo",
     "claude code .   →   /threat-model\n\nOr commit .github/agents/ for Copilot\nOr run teams-app/bot.js for Teams"),
    ("3  Automate",
     "cp .github/workflows/*.yml\n   your-repo/.github/workflows/\n\nAdd ANTHROPIC_API_KEY secret.\nEvery PR → automatic security review."),
]

for i, (title, body) in enumerate(steps):
    x = Inches(0.4 + 4.17 * i)
    rect(s, x, Inches(1.6), Inches(4.0), Inches(4.95), RGBColor(0x16,0x2A,0x3D))
    rect(s, x, Inches(1.6), Inches(4.0), Inches(0.46), AMBER)
    txt(s, title, x + Inches(0.14), Inches(1.65), Inches(3.72), Inches(0.36),
        size=15, bold=True, color=NAVY)
    rect(s, x + Inches(0.12), Inches(2.16), Inches(3.76), Inches(4.29), NAVY)
    txt(s, body, x + Inches(0.24), Inches(2.28), Inches(3.52), Inches(4.07),
        size=12, color=GREEN)

txt(s, "github.com/h0p3sf4ll/threatlint  ·  Read-only  ·  Evidence-gated  ·  Zero infrastructure",
    Inches(1), Inches(6.85), Inches(11.33), Inches(0.42),
    size=13, color=GREY, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "threatlint-executive-overview.pptx")
prs.save(OUT)
print(f"Saved: {OUT}  ({len(prs.slides)} slides)")
